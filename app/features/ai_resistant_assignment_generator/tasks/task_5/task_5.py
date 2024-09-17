import streamlit as st
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.vectorstores import Chroma
import os
import tempfile
import uuid
import pandas as pd
from pptx import Presentation
import docx

class DocumentProcessor:
    """
    This class handles the processing of various document types, including PDFs, DOCX, PPT, 
    and plain text. The processed content is stored as a list of pages/chunks for further use.
    """

    def __init__(self):
        # Initialize an empty list to store pages from all documents
        self.pages = []

    def ingest_documents(self):
        """
        Orchestrates the document ingestion process. Handles different file types (PDF, DOCX, PPT).
        """
        # Call the handler for PDF uploads
        self.handle_pdf_upload()

        # Call the handler for DOCX uploads
        self.handle_docx_upload()

        # Call the handler for PPT uploads
        self.handle_ppt_upload()

        # Call the handler for plain text input
        self.handle_text_input()

    def handle_pdf_upload(self):
        """
        Handles PDF uploads.
        """
        uploaded_files = st.file_uploader("Upload your PDFs", type=['pdf'], accept_multiple_files=True)
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                try:
                    unique_id = uuid.uuid4().hex
                    temp_file_name = f"{uploaded_file.name}_{unique_id}.pdf"
                    temp_file_path = os.path.join(tempfile.gettempdir(), temp_file_name)

                    with open(temp_file_path, 'wb') as f:
                        f.write(uploaded_file.getvalue())

                    loader = PyPDFLoader(temp_file_path)
                    document_pages = loader.load()

                    # Debug: Check the content of the extracted pages
                    for i, page in enumerate(document_pages):
                        st.write(f"Processing page {i+1}:")
                        st.write(page.page_content)  # Ensure page content is correctly extracted
                    
                    # Add the extracted pages to the `pages` list
                    self.pages.extend(document_pages)
                except Exception as e:
                    st.error(f"Error processing file {uploaded_file.name}: {e}")
                finally:
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)

    def handle_docx_upload(self):
        """
        Handles DOCX uploads.
        """
        uploaded_files = st.file_uploader("Upload your DOCX files", type=['docx'], accept_multiple_files=True)
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                try:
                    doc = docx.Document(uploaded_file)
                    for para in doc.paragraphs:
                        self.pages.append(para.text)
                except Exception as e:
                    st.error(f"Error processing DOCX file {uploaded_file.name}: {e}")

    def handle_ppt_upload(self):
        """
        Handles PPT uploads.
        """
        uploaded_file = st.file_uploader("Upload your PPT files", type=['pptx'])
        
        if uploaded_file:
            try:
                unique_id = uuid.uuid4().hex
                temp_file_name = f"{uploaded_file.name}_{unique_id}.pptx"
                temp_file_path = os.path.join(tempfile.gettempdir(), temp_file_name)

                with open(temp_file_path, 'wb') as f:
                    f.write(uploaded_file.getvalue())

                prs = Presentation(temp_file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            self.pages.append(shape.text)
            except Exception as e:
                st.error(f"Error processing PPT file {uploaded_file.name}: {e}")
            finally:
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)

    def handle_text_input(self):
        """
        Handles plain text input.
        """
        text_input = st.text_area("Enter your text here")
        if text_input:
            self.pages.append(text_input.splitlines())

class ChromaCollectionCreator:
    """
    This class handles the creation of a Chroma collection from documents processed by the DocumentProcessor.
    """

    def __init__(self, processor, embed_model):
        self.processor = processor  # Instance of DocumentProcessor
        self.embed_model = embed_model  # Embedding model for Chroma
        self.db = None
        self.persist_directory = r"C:\path\to\chroma_db"  # Update with correct path

        # Ensure persist directory exists
        if not os.path.exists(self.persist_directory):
            os.makedirs(self.persist_directory)

    def create_chroma_collection(self):
        """
        Task: Create a Chroma collection from the documents processed by the DocumentProcessor instance.
        """
        if len(self.processor.pages) == 0:
            st.error("No documents found!", icon="🚨")
            return

        documents = [Document(page_content=str(page)) for page in self.processor.pages]

        # Split documents into chunks for embedding
        text_splitter = CharacterTextSplitter(separator='\n', chunk_size=1000, chunk_overlap=200)
        texts = text_splitter.split_documents(documents)
        if texts:
            st.success(f"Successfully split pages into {len(texts)} documents!", icon="✅")
        else:
            st.error("Failed to split documents into chunks!", icon="🚨")
            return

        # Create the Chroma collection
        try:
            self.db = Chroma.from_documents(documents=texts, embedding=self.embed_model, persist_directory=self.persist_directory)
            st.success("Successfully created Chroma Collection!", icon="✅")
        except Exception as e:
            st.error(f"Failed to create Chroma Collection: {str(e)}", icon="🚨")

    def as_retriever(self):
        if self.db:
            return self.db.as_retriever()
        else:
            raise ValueError("Chroma collection is not initialized.")

    def query_chroma_collection(self, query):
        """
        Queries the Chroma collection for documents similar to the query.
        :param query: The query string to search for in the Chroma collection.
        :return: The first matching document from the collection with similarity score.
        """
        if self.db:
            retriever = self.as_retriever()
            docs = retriever.get_relevant_documents(query)
            if docs:
                return docs[0]  # Returning the first matching document
            else:
                st.error("No matching documents found!", icon="🚨")
                return None
        else:
            st.error("Chroma Collection has not been created!", icon="🚨")


# Main Streamlit app
if __name__ == "__main__":

    st.header("AI-Resistant Assignment Generator")

    # Initialize DocumentProcessor
    processor = DocumentProcessor()

    # Create a form for document ingestion
    with st.form("document_ingestion"):
        st.write("Upload your files or enter text:")

        # Call the ingestion methods
        processor.ingest_documents()

        # Submit button
        submitted = st.form_submit_button("Submit")

        # Check if the form is submitted and process the documents
        if submitted:
            st.success(f"Total pages processed: {len(processor.pages)}")
            for i, page in enumerate(processor.pages):
                st.write(f"Page {i+1}:")
                st.write(page)
