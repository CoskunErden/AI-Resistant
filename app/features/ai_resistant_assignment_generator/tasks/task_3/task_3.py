import streamlit as st
from langchain_community.document_loaders import PyPDFLoader
import os
import tempfile
import uuid
import pandas as pd
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
from bs4 import BeautifulSoup
import requests
from io import StringIO
from googleapiclient.discovery import build
from google.oauth2 import service_account

class DocumentProcessor:
    """
    This class handles the processing of various document types, including PDFs, text, YouTube links,
    Google Sheets/CSV, Google Slides/PPT, web pages, and notes. The processed content is stored as 
    a list of pages/chunks for further use.
    """

    def __init__(self):
        # Initialize an empty list to store pages from all documents
        self.pages = []

    def ingest_documents(self):
        """
        Orchestrates the document ingestion process by calling different handler methods based on the 
        type of input. After processing, it displays the total number of pages/chunks processed.
        """
        # Call the handler for PDF uploads
        self.handle_pdf_upload()
        
        # Call the handler for plain text input
        self.handle_text_input()

        # Call the handler for YouTube links
        self.handle_youtube_input()

        # Call the handler for Google Sheets/CSV uploads
        self.handle_google_sheet_csv_upload()

        # Call the handler for Google Slides/PPT uploads
        self.handle_google_slides_ppt_upload()

        # Call the handler for web page URLs
        self.handle_web_page_input()

        # Call the handler for notes input
        self.handle_notes_input()

        # After all inputs are processed, display the total number of pages/chunks processed
        st.write(f"Total pages processed: {len(self.pages)}")

    def handle_pdf_upload(self):
        """
        Handles the upload and processing of PDF files. Each PDF file is saved temporarily, processed 
        to extract its pages, and then the temporary file is deleted. The extracted pages are added 
        to the `pages` list.
        """
        # Create a file uploader widget for PDFs
        uploaded_files = st.file_uploader(
            "Upload your PDFs",
            type=['pdf'],
            accept_multiple_files=True
        )
        
        if uploaded_files:
            for uploaded_file in uploaded_files:
                try:
                    # Generate a unique identifier to append to the file's original name
                    unique_id = uuid.uuid4().hex
                    original_name, file_extension = os.path.splitext(uploaded_file.name)
                    temp_file_name = f"{original_name}_{unique_id}{file_extension}"
                    temp_file_path = os.path.join(tempfile.gettempdir(), temp_file_name)

                    # Write the uploaded PDF to a temporary file
                    with open(temp_file_path, 'wb') as f:
                        f.write(uploaded_file.getvalue())

                    # Load and process the temporary PDF file using PyPDFLoader
                    loader = PyPDFLoader(temp_file_path)
                    document_pages = loader.load()

                    # Add the extracted pages to the `pages` list
                    self.pages.extend(document_pages)
                except Exception as e:
                    # Handle any errors during processing
                    st.error(f"Error processing file {uploaded_file.name}: {e}")
                finally:
                    # Clean up by deleting the temporary file after processing
                    if os.path.exists(temp_file_path):
                        os.unlink(temp_file_path)

    def handle_text_input(self):
        """
        Handles the input of plain text. The text is split into lines and added to the `pages` list.
        """
        # Create a text area widget for user input
        text_input = st.text_area("Enter your text here")
        
        if text_input:
            # Split the input text into lines and add to the `pages` list
            self.pages.append(text_input.splitlines())

    def handle_youtube_input(self):
        """
        Handles the input of YouTube video links. The transcript is extracted using the YouTube 
        Transcript API, split into lines, and added to the `pages` list.
        """
        # Create a text input widget for YouTube links
        youtube_url = st.text_input("Enter a YouTube link")
        
        if youtube_url:
            try:
                # Extract the video ID from the YouTube URL
                video_id = youtube_url.split("v=")[-1]

                # Retrieve the transcript of the YouTube video using the YouTube Transcript API
                transcript = YouTubeTranscriptApi.get_transcript(video_id)

                # Join the transcript into a single string and split it into lines
                transcript_text = " ".join([t['text'] for t in transcript])
                self.pages.append(transcript_text.splitlines())
            except Exception as e:
                # Handle any errors during processing
                st.error(f"Error processing YouTube link: {e}")

    def handle_google_sheet_csv_upload(self):
        """
        Handles the upload and processing of Google Sheets or CSV files. The CSV file is loaded 
        into a Pandas DataFrame, converted to a string, split into lines, and added to the `pages` list.
        """
        # Create a file uploader widget for Google Sheets/CSV files
        uploaded_file = st.file_uploader("Upload a Google Sheet/CSV", type=['csv'])
        
        if uploaded_file:
            try:
                # Load the CSV file into a Pandas DataFrame
                df = pd.read_csv(uploaded_file)

                # Convert the DataFrame to a string and split it into lines
                self.pages.append(df.to_string().splitlines())
            except Exception as e:
                # Handle any errors during processing
                st.error(f"Error processing file: {e}")

    def handle_google_slides_ppt_upload(self):
        """
        Handles the upload and processing of Google Slides or PPT files. The slides are processed to 
        extract text, which is split into lines and added to the `pages` list.
        """
        # Create a file uploader widget for Google Slides/PPT files
        uploaded_file = st.file_uploader("Upload Google Slides/PPT", type=['pptx'])
        
        if uploaded_file:
            try:
                # Generate a unique identifier to append to the file's original name
                unique_id = uuid.uuid4().hex
                original_name, file_extension = os.path.splitext(uploaded_file.name)
                temp_file_name = f"{original_name}_{unique_id}{file_extension}"
                temp_file_path = os.path.join(tempfile.gettempdir(), temp_file_name)

                # Write the uploaded PPT file to a temporary file
                with open(temp_file_path, 'wb') as f:
                    f.write(uploaded_file.getvalue())

                # Load the presentation and extract text from each slide
                prs = Presentation(temp_file_path)
                for slide in prs.slides:
                    slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                    self.pages.append(slide_text.splitlines())
            except Exception as e:
                # Handle any errors during processing
                st.error(f"Error processing Google Slides/PPT: {e}")
            finally:
                # Clean up by deleting the temporary file after processing
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)

    def handle_web_page_input(self):
        """
        Handles the input of web page URLs. The page content is retrieved, text is extracted using 
        BeautifulSoup, split into lines, and added to the `pages` list.
        """
        # Create a text input widget for web page URLs
        web_url = st.text_input("Enter a web page URL")
        
        if web_url:
            try:
                # Fetch the web page content using requests
                response = requests.get(web_url)

                # Parse the page content with BeautifulSoup to extract text
                soup = BeautifulSoup(response.content, 'html.parser')
                text = soup.get_text()

                # Split the extracted text into lines and add to the `pages` list
                self.pages.append(text.splitlines())
            except Exception as e:
                # Handle any errors during processing
                st.error(f"Error processing web page: {e}")

    def handle_notes_input(self):
        """
        Handles the input of notes. The notes are split into lines and added to the `pages` list.
        """
        # Create a text area widget for notes input
        notes_input = st.text_area("Enter your notes here (Markdown supported)")
        
        if notes_input:
            # Split the notes into lines and add to the `pages` list
            self.pages.append(notes_input.splitlines())

# Run the Streamlit app if this script is executed directly
if __name__ == "__main__":
    processor = DocumentProcessor()
    processor.ingest_documents()
